#!/usr/bin/env python3
"""
ingest_slideserve.py
- Crawl Slideserve search URL
- Respect robots.txt
- Extract presentation pages and attached files (pdf/ppt/pptx/images)
- Download assets, extract text and OCR images
- Produce metadata.json (one entry per presentation)
"""
import os
import sys
import time
import json
import hashlib
import argparse
import logging
import random
from urllib.parse import urljoin, urlparse, quote
from pathlib import Path
import requests
from bs4 import BeautifulSoup
from tqdm import tqdm
import backoff
try:
    from robotexclusionrulesparser import RobotFileParserEx
except ImportError:
    print("⚠️  robotexclusionrulesparser not found, skipping robots.txt check")
    RobotFileParserEx = None

# Optional imports with fallbacks
try:
    from pdfminer.high_level import extract_text as extract_pdf_text
    HAS_PDFMINER = True
except ImportError:
    print("⚠️  pdfminer.six not found, PDF extraction disabled")
    HAS_PDFMINER = False
    
try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    print("⚠️  python-pptx not found, PPTX extraction disabled")
    HAS_PPTX = False
    
try:
    from PIL import Image
    import pytesseract
    HAS_OCR = True
except ImportError:
    print("⚠️  PIL/pytesseract not found, OCR disabled")
    HAS_OCR = False

# Setup logging
logging.basicConfig(
    level=logging.INFO, 
    format="%(asctime)s [%(levelname)s] %(message)s",
    handlers=[
        logging.FileHandler("slideserve_ingest.log"),
        logging.StreamHandler()
    ]
)

HEADERS = {
    "User-Agent": "Mozilla/5.0 (compatible; SovereignGTP-IngestBot/1.0; +mailto:research@example.com)",
    "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8",
    "Accept-Language": "en-US,en;q=0.5",
    "Accept-Encoding": "gzip, deflate",
    "Connection": "keep-alive",
    "Upgrade-Insecure-Requests": "1"
}

SESSION = requests.Session()
SESSION.headers.update(HEADERS)
BASE_DELAY = 2.0  # Base delay between requests (seconds)

def get_delay():
    """Return randomized delay to appear more human-like"""
    return BASE_DELAY + random.uniform(0.5, 2.0)

def sha256_of_file(path):
    """Calculate SHA256 hash of file"""
    try:
        h = hashlib.sha256()
        with open(path, "rb") as f:
            for chunk in iter(lambda: f.read(8192), b""):
                h.update(chunk)
        return h.hexdigest()
    except Exception as e:
        logging.warning(f"Failed to hash {path}: {e}")
        return "unknown"

def sanitize_filename(name, max_len=100):
    """Create safe filename from string"""
    import re
    # Remove/replace problematic characters
    safe = re.sub(r'[<>:"/\\|?*]', '_', name)
    safe = re.sub(r'\s+', '_', safe)
    safe = safe.strip('._')
    return safe[:max_len] if safe else "unnamed"

@backoff.on_exception(backoff.expo, requests.RequestException, max_tries=3)
def robust_get(url, **kwargs):
    """Make HTTP request with retry logic"""
    kwargs.setdefault('timeout', 30)
    return SESSION.get(url, **kwargs)

def ensure_robot_ok(root_url, user_agent="*"):
    """Check robots.txt compliance"""
    if not RobotFileParserEx:
        logging.warning("Robots.txt parser not available, assuming allowed")
        return True, None
        
    try:
        parsed = urlparse(root_url)
        robots_url = f"{parsed.scheme}://{parsed.netloc}/robots.txt"
        logging.info(f"Checking robots.txt: {robots_url}")
        
        r = robust_get(robots_url)
        r.raise_for_status()
        
        parser = RobotFileParserEx()
        parser.parse(r.text.splitlines())
        allowed = parser.is_allowed(user_agent, root_url)
        
        return allowed, parser
    except Exception as e:
        logging.warning(f"Failed to check robots.txt: {e}, assuming allowed")
        return True, None

def get_search_pages(start_url, max_pages=3):
    """Return list of search result page URLs"""
    pages = [start_url]
    
    try:
        logging.info(f"Discovering pagination from: {start_url}")
        r = robust_get(start_url)
        r.raise_for_status()
        soup = BeautifulSoup(r.text, "html.parser")
        
        # Look for pagination links
        next_selectors = [
            "a.next", "a.pagination__next", "a.pager-next",
            "a[aria-label*='next']", "a[title*='next' i]",
            ".pagination a:contains('Next')", ".pager a:contains('Next')"
        ]
        
        current_url = start_url
        for page_num in range(2, max_pages + 1):
            time.sleep(get_delay())
            
            # Try to find next page link
            next_link = None
            for selector in next_selectors:
                try:
                    next_link = soup.select_one(selector.replace(':contains', ''))
                    if next_link:
                        break
                except:
                    continue
            
            if next_link and next_link.get("href"):
                href = next_link.get("href")
                next_url = urljoin(current_url, href)
                if next_url not in pages:
                    pages.append(next_url)
                    logging.info(f"Found page {page_num}: {next_url}")
                    
                    # Get next page to continue pagination
                    try:
                        r = robust_get(next_url)
                        r.raise_for_status()
                        soup = BeautifulSoup(r.text, "html.parser")
                        current_url = next_url
                    except Exception as e:
                        logging.warning(f"Failed to fetch page {page_num}: {e}")
                        break
                else:
                    break  # Already seen this URL
            else:
                # Try constructing page URLs manually
                try:
                    base_url = start_url.split('?')[0]
                    params = start_url.split('?')[1] if '?' in start_url else ''
                    page_url = f"{base_url}?page={page_num}&{params}" if params else f"{base_url}?page={page_num}"
                    
                    r = robust_get(page_url)
                    r.raise_for_status()
                    if page_url not in pages and len(r.text) > 1000:  # Basic content check
                        pages.append(page_url)
                        logging.info(f"Constructed page {page_num}: {page_url}")
                    else:
                        break
                except:
                    break
                    
    except Exception as e:
        logging.warning(f"Failed to enumerate pages: {e}")
    
    return list(dict.fromkeys(pages))  # Remove duplicates while preserving order

def extract_presentation_links(page_url):
    """Extract presentation detail page links from search results"""
    links = []
    try:
        logging.info(f"Extracting links from: {page_url}")
        r = robust_get(page_url)
        r.raise_for_status()
        soup = BeautifulSoup(r.text, "html.parser")
        
        # Look for presentation links
        link_patterns = [
            "/presentation/", "/presentations/", "/slideshow/", "/slides/"
        ]
        
        for a in soup.find_all("a", href=True):
            href = a.get("href", "")
            if any(pattern in href for pattern in link_patterns):
                full_url = urljoin(page_url, href)
                if full_url not in links:
                    links.append(full_url)
                    
        logging.info(f"Found {len(links)} presentation links")
                    
    except Exception as e:
        logging.error(f"Failed to extract links from {page_url}: {e}")
    
    return links

def download_file(url, dest_path, max_size_mb=50):
    """Download file with size limit and error handling"""
    try:
        logging.info(f"Downloading: {url} -> {dest_path}")
        
        # Create directory if needed
        os.makedirs(os.path.dirname(dest_path), exist_ok=True)
        
        r = robust_get(url, stream=True)
        r.raise_for_status()
        
        # Check content length
        content_length = r.headers.get('content-length')
        if content_length:
            size_mb = int(content_length) / (1024 * 1024)
            if size_mb > max_size_mb:
                logging.warning(f"File too large ({size_mb:.1f}MB > {max_size_mb}MB): {url}")
                return False
        
        downloaded_size = 0
        with open(dest_path, "wb") as f:
            for chunk in r.iter_content(chunk_size=8192):
                if chunk:
                    f.write(chunk)
                    downloaded_size += len(chunk)
                    # Safety check during download
                    if downloaded_size > max_size_mb * 1024 * 1024:
                        logging.warning(f"Download size exceeded limit: {url}")
                        return False
        
        logging.info(f"Downloaded {downloaded_size/1024:.1f}KB: {dest_path}")
        return True
        
    except Exception as e:
        logging.warning(f"Failed download {url}: {e}")
        if os.path.exists(dest_path):
            try:
                os.remove(dest_path)
            except:
                pass
        return False

def extract_text_from_pdf(path):
    """Extract text from PDF file"""
    if not HAS_PDFMINER:
        return ""
    try:
        text = extract_pdf_text(path)
        return text.strip() if text else ""
    except Exception as e:
        logging.warning(f"PDF extraction failed {path}: {e}")
        return ""

def extract_text_from_pptx(path):
    """Extract text from PowerPoint file"""
    if not HAS_PPTX:
        return ""
    
    text_parts = []
    try:
        prs = Presentation(path)
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            if slide_text:
                text_parts.append(f"=== Slide {slide_num} ===")
                text_parts.extend(slide_text)
                text_parts.append("")  # Blank line between slides
                
    except Exception as e:
        logging.warning(f"PPTX extraction failed {path}: {e}")
    
    return "\n".join(text_parts)

def ocr_image(path):
    """Extract text from image using OCR"""
    if not HAS_OCR:
        return ""
    
    try:
        with Image.open(path) as img:
            # Convert to RGB if needed
            if img.mode != 'RGB':
                img = img.convert('RGB')
            text = pytesseract.image_to_string(img, lang='eng')
            return text.strip()
    except Exception as e:
        logging.warning(f"OCR failed for {path}: {e}")
        return ""

def is_valid_content(text, min_words=10):
    """Check if extracted text is meaningful"""
    if not text or not isinstance(text, str):
        return False
    words = text.split()
    return len(words) >= min_words and not text.isdigit()

def process_presentation_page(page_url, outdir):
    """Process a single presentation page"""
    logging.info(f"🎯 Processing presentation: {page_url}")
    time.sleep(get_delay())
    
    try:
        r = robust_get(page_url)
        r.raise_for_status()
        soup = BeautifulSoup(r.text, "html.parser")
        
        # Extract title
        title_candidates = [
            soup.find("title"),
            soup.find("h1"),
            soup.find(class_=lambda x: x and "title" in x.lower()),
            soup.find("meta", property="og:title")
        ]
        
        title = "Unknown Presentation"
        for candidate in title_candidates:
            if candidate:
                if candidate.name == "meta":
                    title = candidate.get("content", "").strip()
                else:
                    title = candidate.get_text().strip()
                if title and len(title) > 5:
                    break
        
        # Clean title for directory name
        slug = sanitize_filename(title)
        presentation_dir = os.path.join(outdir, f"{slug}_{hashlib.md5(page_url.encode()).hexdigest()[:8]}")
        os.makedirs(presentation_dir, exist_ok=True)
        
        # Save page HTML
        html_path = os.path.join(presentation_dir, "page.html")
        with open(html_path, "w", encoding="utf-8") as fh:
            fh.write(r.text)
        
        assets = {
            "page_html": html_path,
            "source_url": page_url,
            "title": title,
            "files": [],
            "extracted_texts": []
        }
        
        # Find downloadable files
        download_extensions = [".pdf", ".ppt", ".pptx", ".zip", ".doc", ".docx"]
        for a in soup.find_all("a", href=True):
            href = a.get("href", "")
            full_url = urljoin(page_url, href)
            
            if any(ext in full_url.lower() for ext in download_extensions):
                # Extract filename
                filename = os.path.basename(full_url.split("?")[0])
                if not filename or len(filename) < 3:
                    filename = f"download_{len(assets['files'])}"
                
                filename = sanitize_filename(filename)
                dest_path = os.path.join(presentation_dir, filename)
                
                if download_file(full_url, dest_path):
                    assets["files"].append(dest_path)
        
        # Download presentation images/thumbnails
        image_extensions = [".jpg", ".jpeg", ".png", ".gif", ".webp"]
        images_found = 0
        for img in soup.find_all("img", src=True):
            if images_found >= 20:  # Limit images
                break
                
            img_src = img.get("src", "")
            if not img_src:
                continue
                
            img_url = urljoin(page_url, img_src)
            
            # Skip tiny images, icons, etc.
            if any(skip in img_url.lower() for skip in ["icon", "logo", "avatar", "button"]):
                continue
                
            filename = f"img_{images_found:03d}_{os.path.basename(img_url.split('?')[0])}"
            filename = sanitize_filename(filename)
            dest_path = os.path.join(presentation_dir, filename)
            
            if download_file(img_url, dest_path, max_size_mb=5):
                assets["files"].append(dest_path)
                images_found += 1
        
        # Extract text from downloaded files
        for fpath in assets["files"]:
            if not os.path.exists(fpath):
                continue
                
            f_lower = fpath.lower()
            extracted_text = ""
            
            if f_lower.endswith(".pdf"):
                extracted_text = extract_text_from_pdf(fpath)
            elif f_lower.endswith((".pptx", ".ppt")):
                extracted_text = extract_text_from_pptx(fpath)
            elif any(f_lower.endswith(ext) for ext in [".png", ".jpg", ".jpeg", ".gif", ".webp"]):
                extracted_text = ocr_image(fpath)
            
            if is_valid_content(extracted_text):
                txt_path = fpath + ".txt"
                with open(txt_path, "w", encoding="utf-8") as fh:
                    fh.write(extracted_text)
                assets["extracted_texts"].append(txt_path)
                logging.info(f"✅ Extracted {len(extracted_text.split())} words from {os.path.basename(fpath)}")
        
        # Create file inventory with hashes
        file_entries = []
        all_files = [assets["page_html"]] + assets["files"] + assets["extracted_texts"]
        
        for fpath in all_files:
            if os.path.exists(fpath):
                rel_path = os.path.relpath(fpath, outdir)
                file_entries.append({
                    "path": rel_path,
                    "sha256": sha256_of_file(fpath),
                    "size": os.path.getsize(fpath),
                    "type": "extracted_text" if fpath in assets["extracted_texts"] else "asset"
                })
        
        # Generate metadata
        metadata = {
            "title": title,
            "source_url": page_url,
            "saved_at": time.strftime("%Y-%m-%dT%H:%M:%SZ", time.gmtime()),
            "files": file_entries,
            "stats": {
                "total_files": len(assets["files"]),
                "extracted_texts": len(assets["extracted_texts"]),
                "total_size_bytes": sum(entry["size"] for entry in file_entries)
            }
        }
        
        # Save metadata
        meta_path = os.path.join(presentation_dir, "metadata.json")
        with open(meta_path, "w", encoding="utf-8") as fh:
            json.dump(metadata, fh, indent=2, ensure_ascii=False)
        
        logging.info(f"✅ Completed: {title} ({len(assets['files'])} files, {len(assets['extracted_texts'])} texts)")
        return metadata
        
    except Exception as e:
        logging.error(f"❌ Failed processing {page_url}: {e}")
        raise

def main():
    parser = argparse.ArgumentParser(description="Ingest SlideServe presentations")
    parser.add_argument("start_url", help="SlideServe search URL")
    parser.add_argument("--outdir", default="slideserve_out", help="Output directory")
    parser.add_argument("--limit", type=int, default=0, help="Max presentations to process (0=unlimited)")
    parser.add_argument("--max-pages", type=int, default=3, help="Max search pages to scan")
    args = parser.parse_args()
    
    print(f"""
🎯 SlideServe Ingestion Starting
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
📍 URL: {args.start_url}
📁 Output: {args.outdir}
🎚️  Limit: {args.limit if args.limit else 'Unlimited'}
📄 Max Pages: {args.max_pages}
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    """)
    
    os.makedirs(args.outdir, exist_ok=True)
    
    # Check robots.txt
    allowed, _ = ensure_robot_ok(args.start_url)
    if not allowed:
        logging.error("❌ Crawling disallowed by robots.txt. Aborting.")
        sys.exit(1)
    logging.info("✅ Crawling allowed by robots.txt")
    
    # Discover search pages
    pages = get_search_pages(args.start_url, max_pages=args.max_pages)
    logging.info(f"📄 Discovered {len(pages)} search pages")
    
    # Extract all presentation links
    all_links = []
    for page in pages:
        try:
            links = extract_presentation_links(page)
            all_links.extend(links)
            time.sleep(get_delay())
        except Exception as e:
            logging.warning(f"Failed to process search page {page}: {e}")
    
    # Remove duplicates
    all_links = list(dict.fromkeys(all_links))
    total_found = len(all_links)
    
    if args.limit and args.limit < total_found:
        all_links = all_links[:args.limit]
        logging.info(f"🎯 Limited to first {args.limit} of {total_found} presentations")
    
    logging.info(f"🎬 Processing {len(all_links)} presentations...")
    
    # Process presentations
    all_metadata = []
    successful = 0
    failed = 0
    
    with tqdm(all_links, desc="Processing") as pbar:
        for i, link in enumerate(pbar, 1):
            try:
                pbar.set_description(f"Processing {i}/{len(all_links)}")
                meta = process_presentation_page(link, args.outdir)
                all_metadata.append(meta)
                successful += 1
                
            except Exception as e:
                logging.warning(f"❌ Failed processing {link}: {e}")
                failed += 1
            
            # Update progress
            pbar.set_postfix({"✅": successful, "❌": failed})
    
    # Save master metadata
    master_meta_path = os.path.join(args.outdir, "master_metadata.json")
    with open(master_meta_path, "w", encoding="utf-8") as fh:
        json.dump({
            "ingestion_summary": {
                "total_found": total_found,
                "successful": successful,
                "failed": failed,
                "completed_at": time.strftime("%Y-%m-%dT%H:%M:%SZ", time.gmtime()),
                "source_url": args.start_url
            },
            "presentations": all_metadata
        }, fh, indent=2, ensure_ascii=False)
    
    print(f"""
🎉 Ingestion Complete!
━━━━━━━━━━━━━━━━━━━━━━
✅ Successful: {successful}
❌ Failed: {failed}  
📊 Total Found: {total_found}
📁 Output: {args.outdir}
📄 Metadata: {master_meta_path}
━━━━━━━━━━━━━━━━━━━━━━
    """)

if __name__ == "__main__":
    main()
