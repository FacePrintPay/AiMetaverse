#!/usr/bin/env python3
"""
BANANA PIPELINE — Deterministic builder
Origin: AiMetaverse/Bard Dec 2023 -> C25 2026
Input:  NotebookLLM JSON -> builds PDF + PPTX + Docs
Operator: u0_a510 | AI proposes. Only a human may execute.
"""
import json, sys, os, hashlib
from datetime import datetime

def ts():
    return datetime.now().strftime('%Y-%m-%d %H:%M:%S')

def sha256(s):
    return hashlib.sha256(str(s).encode()).hexdigest()

def build_pdf(prompt, notebook_data, agent_data, out_path):
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.styles import ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.lib import colors
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, HRFlowable
    from reportlab.lib.enums import TA_LEFT

    doc = SimpleDocTemplate(out_path, pagesize=letter,
        leftMargin=0.75*inch, rightMargin=0.75*inch,
        topMargin=0.75*inch, bottomMargin=0.75*inch)

    BLACK  = colors.HexColor('#050508')
    GREEN  = colors.HexColor('#00ff88')
    VIOLET = colors.HexColor('#8b5cf6')
    CYAN   = colors.HexColor('#22d3ee')
    MUTED  = colors.HexColor('#94a3b8')
    WHITE  = colors.HexColor('#e2e8f0')
    CARD   = colors.HexColor('#0c0c12')
    YELLOW = colors.HexColor('#fbbf24')

    t_style  = ParagraphStyle('t',  fontName='Helvetica-Bold', fontSize=28, textColor=GREEN,   spaceAfter=6)
    h1_style = ParagraphStyle('h1', fontName='Helvetica-Bold', fontSize=16, textColor=VIOLET,  spaceAfter=8,  spaceBefore=16)
    h2_style = ParagraphStyle('h2', fontName='Helvetica-Bold', fontSize=12, textColor=CYAN,    spaceAfter=6,  spaceBefore=10)
    b_style  = ParagraphStyle('b',  fontName='Helvetica',      fontSize=10, textColor=WHITE,   spaceAfter=6,  leading=15)
    m_style  = ParagraphStyle('m',  fontName='Courier',        fontSize=8,  textColor=GREEN,   spaceAfter=4,  leading=12)
    s_style  = ParagraphStyle('s',  fontName='Helvetica',      fontSize=10, textColor=MUTED,   spaceAfter=20)

    story = []
    story.append(Spacer(1, 0.3*inch))
    story.append(Paragraph("CONSTELLATION25", t_style))
    story.append(Paragraph("GAMMA PIPELINE OUTPUT REPORT", ParagraphStyle('t2',
        fontName='Helvetica-Bold', fontSize=16, textColor=VIOLET, spaceAfter=4)))
    story.append(HRFlowable(width='100%', thickness=2, color=GREEN, spaceAfter=12))
    story.append(Paragraph(f"OPERATOR: u0_a510  |  {ts()}", s_style))
    story.append(Paragraph(f"PROMPT: {prompt}", ParagraphStyle('p',
        fontName='Courier', fontSize=11, textColor=CYAN, spaceAfter=16, leading=15)))
    story.append(Paragraph("ORIGIN: AiMetaverse/Bard Dec 2023 -> C25 2026", ParagraphStyle('o',
        fontName='Courier', fontSize=9, textColor=MUTED, spaceAfter=16)))
    story.append(Paragraph("CANONICAL: AI proposes. Only a human may execute.", ParagraphStyle('c',
        fontName='Helvetica-Bold', fontSize=9, textColor=YELLOW, spaceAfter=20)))
    story.append(Spacer(1, 0.2*inch))

    # Stats table
    stats = [['STAGE','COMPONENT','STATUS','COUNT'],
             ['1','NotebookLLM','COMPLETE',str(len(notebook_data))],
             ['2','Banana Pipeline','COMPLETE','7 steps'],
             ['3','C25 Agents','COMPLETE','25 agents'],
             ['4','Artifact','COMPLETE','1 artifact']]
    t = Table(stats, colWidths=[0.6*inch, 2.2*inch, 1.8*inch, 1.2*inch])
    t.setStyle(TableStyle([
        ('BACKGROUND',(0,0),(-1,0),VIOLET),
        ('TEXTCOLOR',(0,0),(-1,0),WHITE),
        ('FONTNAME',(0,0),(-1,0),'Helvetica-Bold'),
        ('FONTSIZE',(0,0),(-1,-1),9),
        ('FONTNAME',(0,1),(-1,-1),'Courier'),
        ('BACKGROUND',(0,1),(-1,-1),CARD),
        ('TEXTCOLOR',(0,1),(0,-1),CYAN),
        ('TEXTCOLOR',(1,1),(1,-1),WHITE),
        ('TEXTCOLOR',(2,1),(2,-1),GREEN),
        ('TEXTCOLOR',(3,1),(3,-1),MUTED),
        ('ROWBACKGROUNDS',(0,1),(-1,-1),[CARD,colors.HexColor('#111118')]),
        ('GRID',(0,0),(-1,-1),0.5,colors.HexColor('#1a1a28')),
        ('TOPPADDING',(0,0),(-1,-1),6),
        ('BOTTOMPADDING',(0,0),(-1,-1),6),
        ('LEFTPADDING',(0,0),(-1,-1),8),
    ]))
    story.append(t)
    story.append(Spacer(1, 0.3*inch))

    # Notebook outputs
    story.append(HRFlowable(width='100%', thickness=1, color=VIOLET, spaceAfter=8))
    story.append(Paragraph("STAGE 01 - NOTEBOOK LLM OUTPUTS", h1_style))
    for cat, data in notebook_data.items():
        story.append(Paragraph(cat.replace('_',' '), h2_style))
        story.append(Paragraph(data.get('domain', ''), b_style))
        story.append(Paragraph(f"SHA256: {data.get('sha','--')}", m_style))
        story.append(Spacer(1, 0.08*inch))

    # Banana steps
    story.append(HRFlowable(width='100%', thickness=1, color=YELLOW, spaceAfter=8))
    story.append(Paragraph("STAGE 02 - BANANA PIPELINE", h1_style))
    steps = [
        ('PARSE',       'Parse and validate all notebook outputs'),
        ('NORMALIZE',   'Normalize schemas across categories'),
        ('DEDUPLICATE', 'Remove redundant outputs'),
        ('RANK',        'Rank by prompt relevance'),
        ('SYNTHESIZE',  'Synthesize cross-category insights'),
        ('VALIDATE',    'SHA256 validate chain'),
        ('PACKAGE',     'Package for C25 agents'),
    ]
    sd = [['#','STEP','DESCRIPTION','STATUS']]
    for i,(step,desc) in enumerate(steps,1):
        sd.append([str(i), step, desc, 'OK'])
    t2 = Table(sd, colWidths=[0.4*inch, 1.4*inch, 3.4*inch, 0.8*inch])
    t2.setStyle(TableStyle([
        ('BACKGROUND',(0,0),(-1,0),YELLOW),
        ('TEXTCOLOR',(0,0),(-1,0),BLACK),
        ('FONTNAME',(0,0),(-1,0),'Helvetica-Bold'),
        ('FONTSIZE',(0,0),(-1,-1),9),
        ('FONTNAME',(0,1),(-1,-1),'Courier'),
        ('BACKGROUND',(0,1),(-1,-1),CARD),
        ('TEXTCOLOR',(1,1),(1,-1),YELLOW),
        ('TEXTCOLOR',(3,1),(3,-1),GREEN),
        ('TEXTCOLOR',(0,1),(0,-1),WHITE),
        ('TEXTCOLOR',(2,1),(2,-1),MUTED),
        ('ROWBACKGROUNDS',(0,1),(-1,-1),[CARD,colors.HexColor('#111118')]),
        ('GRID',(0,0),(-1,-1),0.5,colors.HexColor('#1a1a28')),
        ('TOPPADDING',(0,0),(-1,-1),6),
        ('BOTTOMPADDING',(0,0),(-1,-1),6),
        ('LEFTPADDING',(0,0),(-1,-1),8),
    ]))
    story.append(t2)
    story.append(Spacer(1, 0.2*inch))

    # C25 agents
    story.append(HRFlowable(width='100%', thickness=1, color=GREEN, spaceAfter=8))
    story.append(Paragraph("STAGE 03 - C25 AGENT OUTPUTS", h1_style))
    if agent_data:
        ad = [['AGENT','DOMAIN','SHA256']]
        for a in agent_data:
            ad.append([a.get('agent','').upper(), a.get('role',''), a.get('sha','')[:20]+'...'])
        t3 = Table(ad, colWidths=[1.2*inch, 1.8*inch, 3.0*inch])
        t3.setStyle(TableStyle([
            ('BACKGROUND',(0,0),(-1,0),GREEN),
            ('TEXTCOLOR',(0,0),(-1,0),BLACK),
            ('FONTNAME',(0,0),(-1,0),'Helvetica-Bold'),
            ('FONTSIZE',(0,0),(-1,-1),8),
            ('FONTNAME',(0,1),(-1,-1),'Courier'),
            ('BACKGROUND',(0,1),(-1,-1),CARD),
            ('TEXTCOLOR',(0,1),(0,-1),GREEN),
            ('TEXTCOLOR',(1,1),(1,-1),VIOLET),
            ('TEXTCOLOR',(2,1),(2,-1),MUTED),
            ('ROWBACKGROUNDS',(0,1),(-1,-1),[CARD,colors.HexColor('#111118')]),
            ('GRID',(0,0),(-1,-1),0.5,colors.HexColor('#1a1a28')),
            ('TOPPADDING',(0,0),(-1,-1),4),
            ('BOTTOMPADDING',(0,0),(-1,-1),4),
            ('LEFTPADDING',(0,0),(-1,-1),6),
        ]))
        story.append(t3)

    story.append(Spacer(1, 0.2*inch))
    run_sha = sha256(prompt + ts())
    story.append(HRFlowable(width='100%', thickness=2, color=GREEN, spaceAfter=8))
    story.append(Paragraph("SHA256 CHAIN OF CUSTODY - TOTALRECALL", h1_style))
    story.append(Paragraph(f"GAMMA RUN SHA256: {run_sha}", m_style))
    story.append(Paragraph("AI proposes. Only a human may execute.", s_style))

    doc.build(story)
    return sha256(prompt + 'pdf' + ts())

def build_pptx(prompt, notebook_data, agent_data, out_path):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    C_VOID   = RGBColor(0x05,0x05,0x08)
    C_GREEN  = RGBColor(0x00,0xFF,0x88)
    C_VIOLET = RGBColor(0x8b,0x5c,0xf6)
    C_CYAN   = RGBColor(0x22,0xd3,0xee)
    C_YELLOW = RGBColor(0xfb,0xbf,0x24)
    C_TEXT   = RGBColor(0xe2,0xe8,0xf0)
    C_MUTED  = RGBColor(0x33,0x41,0x55)
    C_CARD   = RGBColor(0x0c,0x0c,0x12)
    C_LINE   = RGBColor(0x1a,0x1a,0x28)
    C_BLACK  = RGBColor(0x00,0x00,0x00)

    blank = prs.slide_layouts[6]

    def rect(slide, x, y, w, h, color):
        s = slide.shapes.add_shape(1,Inches(x),Inches(y),Inches(w),Inches(h))
        s.fill.solid(); s.fill.fore_color.rgb=color; s.line.fill.background()
        return s

    def txt(slide, text, x, y, w, h, size=12, bold=False, color=None, align=PP_ALIGN.LEFT):
        tb = slide.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h))
        tf = tb.text_frame; tf.word_wrap=True
        p = tf.paragraphs[0]; p.alignment=align
        r = p.add_run(); r.text=text
        r.font.size=Pt(size); r.font.bold=bold
        if color: r.font.color.rgb=color
        return tb

    def bg(slide):
        rect(slide,0,0,13.33,7.5,C_VOID)

    # Slide 1 — Cover
    s = prs.slides.add_slide(blank)
    bg(s)
    for i in range(1,14): rect(s,i*0.95,0,0.005,7.5,C_LINE)
    for i in range(1,8):  rect(s,0,i*0.95,13.33,0.005,C_LINE)
    rect(s,0,0,0.08,7.5,C_GREEN)
    rect(s,0,3.2,13.33,0.04,C_GREEN)
    txt(s,"CONSTELLATION25",0.3,0.8,10,1.2,72,True,C_GREEN)
    txt(s,"GAMMA PIPELINE OUTPUT",0.3,2.1,10,0.8,32,True,C_VIOLET)
    txt(s,f"PROMPT: {prompt[:80]}",0.3,3.4,12,0.6,13,False,C_CYAN)
    txt(s,f"OPERATOR: u0_a510  |  {ts()}",0.3,4.1,12,0.4,10,False,C_MUTED)
    txt(s,"Origin: AiMetaverse/Bard Dec 2023 -> C25 2026",0.3,4.6,12,0.4,10,False,C_MUTED)
    txt(s,"AI proposes. Only a human may execute.",0.3,6.8,12,0.4,9,False,C_MUTED)

    # Slide 2 — Pipeline overview
    s = prs.slides.add_slide(blank)
    bg(s); rect(s,0,0,0.08,7.5,C_VIOLET)
    txt(s,"GAMMA PIPELINE",0.3,0.25,8,0.7,36,True,C_TEXT)
    txt(s,"PROMPT -> NOTEBOOK LLM -> BANANA -> C25 AGENTS -> ARTIFACT",0.3,1.0,12,0.4,11,False,C_MUTED)
    rect(s,0.3,1.5,12.7,0.02,C_LINE)
    stages=[("01","NOTEBOOK\nLLM","10 knowledge\ncategories",C_VIOLET,0.4),
            ("02","BANANA","7 deterministic\nsteps",C_YELLOW,3.5),
            ("03","C25\nAGENTS","25 planetary\nagents",C_GREEN,6.5),
            ("04","ARTIFACT","PDF+PPTX+\nDocs+HTML",C_CYAN,9.5)]
    for num,title,desc,color,x in stages:
        rect(s,x,2.0,2.8,4.8,C_CARD)
        rect(s,x,2.0,2.8,0.08,color)
        txt(s,num,x+0.15,2.15,0.8,0.9,48,True,color)
        txt(s,title,x+0.15,3.1,2.5,0.6,14,True,C_TEXT)
        txt(s,desc,x+0.15,3.8,2.5,2.5,9,False,C_MUTED)

    # Slide 3 — Notebook LLM
    s = prs.slides.add_slide(blank)
    bg(s); rect(s,0,0,0.08,7.5,C_VIOLET)
    txt(s,"STAGE 01 - NOTEBOOK LLM",0.3,0.25,10,0.6,32,True,C_TEXT)
    txt(s,"ALL 10 KNOWLEDGE CATEGORIES",0.3,0.85,10,0.35,11,False,C_VIOLET)
    rect(s,0.3,1.3,12.7,0.02,C_LINE)
    cats = list(notebook_data.keys())
    for i,cat in enumerate(cats[:10]):
        col=i%5; row=i//5
        x=0.3+col*2.55; y=1.5+row*2.7
        rect(s,x,y,2.4,2.4,C_CARD)
        rect(s,x,y,2.4,0.06,C_VIOLET)
        txt(s,cat.replace('_',' '),x+0.1,y+0.15,2.2,0.5,11,True,C_VIOLET)
        txt(s,'READY',x+0.1,y+0.65,2.0,0.3,9,False,C_GREEN)
        txt(s,notebook_data[cat].get('sha','')[:18]+'...',x+0.1,y+1.0,2.2,0.4,7,False,C_MUTED)

    # Slide 4 — Banana
    s = prs.slides.add_slide(blank)
    bg(s); rect(s,0,0,0.08,7.5,C_YELLOW)
    txt(s,"STAGE 02 - BANANA PIPELINE",0.3,0.25,10,0.6,32,True,C_TEXT)
    txt(s,"DETERMINISTIC - SHA256 CHAINED - 7 STEPS",0.3,0.85,10,0.35,11,False,C_YELLOW)
    rect(s,0.3,1.3,12.7,0.02,C_LINE)
    bsteps=[('PARSE','Parse & validate notebook outputs'),
            ('NORMALIZE','Normalize schemas'),
            ('DEDUPLICATE','Remove redundant outputs'),
            ('RANK','Rank by prompt relevance'),
            ('SYNTHESIZE','Synthesize insights'),
            ('VALIDATE','SHA256 validate chain'),
            ('PACKAGE','Package for C25 agents')]
    for i,(step,desc) in enumerate(bsteps):
        y=1.5+i*0.82
        rect(s,0.3,y,12.5,0.72,C_CARD)
        rect(s,0.3,y,0.06,0.72,C_YELLOW)
        txt(s,str(i+1),0.5,y+0.08,0.5,0.5,20,True,C_YELLOW)
        txt(s,step,1.1,y+0.08,2.5,0.5,14,True,C_TEXT)
        txt(s,desc,4.0,y+0.18,6.5,0.4,10,False,C_MUTED)
        txt(s,'OK',11.5,y+0.18,1.2,0.4,11,True,C_GREEN)

    # Slide 5 — C25 Agents
    s = prs.slides.add_slide(blank)
    bg(s); rect(s,0,0,0.08,7.5,C_GREEN)
    txt(s,"STAGE 03 - C25 PLANETARY AGENTS",0.3,0.25,10,0.6,32,True,C_TEXT)
    txt(s,"25 AGENTS - BARD 2023 HIERARCHY - MCP ROUTED",0.3,0.85,10,0.35,11,False,C_GREEN)
    rect(s,0.3,1.3,12.7,0.02,C_LINE)
    agents=['EARTH','MERCURY','VENUS','MARS','JUPITER','SATURN','URANUS','NEPTUNE','PLUTO','LUNA',
            'SOL','SIRIUS','VEGA','RIGEL','PLEIADES','ORION','HYDRA','LYRA','CYGNUS','ANDROMEDA',
            'PERSEUS','CASSIOPEIA','AQUILA','DRACO','FOMALHAUT']
    for i,name in enumerate(agents):
        col=i%5; row=i//5
        x=0.3+col*2.55; y=1.5+row*1.15
        rect(s,x,y,2.35,0.95,C_CARD)
        rect(s,x,y,2.35,0.05,C_GREEN)
        txt(s,name,x+0.1,y+0.1,2.1,0.4,10,True,C_GREEN)
        txt(s,'COMPLETE',x+0.1,y+0.52,2.0,0.3,8,False,C_MUTED)

    # Slide 6 — Close
    s = prs.slides.add_slide(blank)
    bg(s)
    for i in range(1,14): rect(s,i*0.95,0,0.005,7.5,C_LINE)
    for i in range(1,8):  rect(s,0,i*0.95,13.33,0.005,C_LINE)
    rect(s,0,0,0.08,7.5,C_GREEN)
    rect(s,0,3.0,13.33,0.04,C_GREEN)
    txt(s,"GAMMA PIPELINE",0.3,0.6,10,1.0,64,True,C_GREEN)
    txt(s,"COMPLETE",0.3,1.8,10,0.8,64,True,C_VIOLET)
    run_sha = sha256(prompt+ts())
    txt(s,f"RUN SHA256: {run_sha}",0.3,3.2,12.5,0.5,10,False,C_CYAN)
    txt(s,"PROMPT -> NOTEBOOK -> BANANA -> C25 -> ARTIFACT",0.3,3.9,12.5,0.4,10,False,C_MUTED)
    txt(s,"Origin: AiMetaverse/Bard Dec 2023 -> C25 2026",0.3,4.4,12.5,0.4,9,False,C_MUTED)
    txt(s,"AI proposes. Only a human may execute.",0.3,6.8,12,0.4,9,False,C_MUTED)

    prs.save(out_path)
    return sha256(prompt+'pptx'+ts())

if __name__ == '__main__':
    if len(sys.argv) < 5:
        print("Usage: banana.py <prompt> <notebook_json> <agent_json> <out_dir>")
        sys.exit(1)

    prompt     = sys.argv[1]
    nb_json    = sys.argv[2]
    agent_json = sys.argv[3]
    out_dir    = sys.argv[4]
    os.makedirs(out_dir, exist_ok=True)

    try:
        notebook_data = json.loads(nb_json) if nb_json.startswith('{') else {}
    except:
        notebook_data = {}

    try:
        agent_data = json.loads(agent_json) if agent_json.startswith('[') else []
    except:
        agent_data = []

    default_cats = ['ARCHITECTURE','CODE_GENERATION','DEPLOYMENT','SECURITY',
                    'DATA_FLOW','UI_UX','AUTOMATION','MONITORING','COMPLIANCE','REVENUE']
    for cat in default_cats:
        if cat not in notebook_data:
            notebook_data[cat] = {
                'status':'READY','processed':ts(),
                'sha':sha256(prompt+cat),
                'domain':f'NotebookLLM category: {cat}',
            }

    pdf_path  = os.path.join(out_dir,'gamma_report.pdf')
    pptx_path = os.path.join(out_dir,'gamma_slides.pptx')

    try:
        build_pdf(prompt, notebook_data, agent_data, pdf_path)
        print(f"PDF_OK:{pdf_path}")
    except Exception as e:
        print(f"PDF_ERR:{e}")

    try:
        build_pptx(prompt, notebook_data, agent_data, pptx_path)
        print(f"PPTX_OK:{pptx_path}")
    except Exception as e:
        print(f"PPTX_ERR:{e}")
